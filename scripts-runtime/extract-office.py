#!/usr/bin/env python3
"""Fallback office file extractor for SAG.

Usage:
    python scripts/extract-office.py --input <file> --output <file>

Reads .xlsx/.xls/.docx/.pptx using pure Python libraries (no COM/Office).
Writes the extracted text as markdown to the output file. The output file
path is also printed on stdout so the Node caller can read it back.
"""
import argparse
import os
import sys

MAX_DETAIL_ROWS = 100
MAX_DETAIL_COLS = 12
DETAIL_KEYWORDS = [
    "明细", "台账", "报表", "记录", "流水", "汇总",
    "details", "ledger", "log", "register", "明细表", "台账表", "汇总表"
]


def _matches_detail_keyword(name: str) -> str | None:
    lower = name.lower()
    for keyword in DETAIL_KEYWORDS:
        if keyword in lower:
            return f"detail-report keyword '{keyword}' matched in '{name}'"
    return None


def read_xlsx(path: str) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    file_name = path.lower()
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"## {ws.title}\n")
        # Skip detail-report sheets: large tables should not be vectorized.
        keyword_match = _matches_detail_keyword(ws.title) or _matches_detail_keyword(file_name)
        if ws.max_row > 200 or ws.max_column > 26 or keyword_match:
            reason = keyword_match if keyword_match else f"{ws.max_row} rows × {ws.max_column} cols"
            out.append(f"*(sheet skipped: {reason} — treated as detail data, not vectorized)*\n")
            continue
        rows: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            rows.append([str(cell) if cell is not None else "" for cell in row])
        if not rows:
            out.append("*(empty sheet)*\n")
            continue
        max_cols = max(len(r) for r in rows)
        for r in rows:
            while len(r) < max_cols:
                r.append("")
        header = rows[0]
        out.append("| " + " | ".join(header) + " |")
        out.append("| " + " | ".join(["---"] * max_cols) + " |")
        for r in rows[1:]:
            out.append("| " + " | ".join(r) + " |")
        out.append("")
    wb.close()
    return "\n".join(out)


def read_xls(path: str) -> str:
    import xlrd

    bk = xlrd.open_workbook(path)
    out: list[str] = []
    for si in range(bk.nsheets):
        sh = bk.sheet_by_index(si)
        out.append(f"## {sh.name}\n")
        if sh.nrows == 0:
            out.append("*(empty sheet)*\n")
            continue
        header = [str(sh.cell_value(0, ci)) for ci in range(sh.ncols)]
        out.append("| " + " | ".join(header) + " |")
        out.append("| " + " | ".join(["---"] * sh.ncols) + " |")
        for ri in range(1, sh.nrows):
            row = [str(sh.cell_value(ri, ci)) for ci in range(sh.ncols)]
            out.append("| " + " | ".join(row) + " |")
        out.append("")
    return "\n".join(out)


def read_docx(path: str) -> str:
    from docx import Document

    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def read_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
    return "\n".join(texts)


def extract(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".xlsx":
        return read_xlsx(path)
    if ext == ".xls":
        return read_xls(path)
    if ext == ".docx":
        return read_docx(path)
    if ext == ".pptx":
        return read_pptx(path)
    raise ValueError(f"unsupported extension: {ext}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Fallback office extractor for SAG")
    parser.add_argument("--input", required=True, help="input file path")
    parser.add_argument("--output", required=True, help="output markdown file path")
    args = parser.parse_args()

    if not os.path.isfile(args.input):
        print(f"input not found: {args.input}", file=sys.stderr)
        sys.exit(10)

    try:
        text = extract(args.input)
    except Exception as exc:
        print(f"extraction failed: {exc}", file=sys.stderr)
        sys.exit(1)

    with open(args.output, "w", encoding="utf-8") as f:
        f.write(text)

    print(args.output)


if __name__ == "__main__":
    main()
